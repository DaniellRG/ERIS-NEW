"""Tool de documentos: lee, extrae y edita archivos de todo tipo.

Formatos soportados:
  - PDF  (.pdf)                          via PyPDF2
  - Word (.docx, .docm)                  via python-docx
  - Excel (.xlsx, .xlsm, .xlsb, .xltx)   via openpyxl
  - PowerPoint (.pptx, .pptm, .potx)     via python-pptx
  - Texto/codigo (txt, md, csv, json, xml, html, py, js, ts, java,
    c, cpp, go, rs, rb, php, sql, sh, bat, ps1, ini, yaml, toml, ...)

Acciones:
  - info / read-info: metadatos del archivo (nombre, tamano, tipo, formato).
  - read / extract / leer: extrae el texto completo (o 'max_chars').
  - summary / resumen: extrae un resumen/preview del contenido.
  - write / save: escribe/sobrescribe un archivo de texto con 'content'.
  - edit / replace: reemplaza en un archivo de texto 'find' por 'replace'
    (si no se pasa 'replace', borra todas las apariciones).
  - append: agrega 'content' al final de un archivo de texto.
  - to_txt / convert: convierte un documento (PDF/Word/Excel/PPT) a .txt/.md.
  - formats / list: lista los formatos soportados.
"""

import os
import re
from datetime import datetime
from pathlib import Path

try:
    from PyPDF2 import PdfReader
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    import docx as _docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import requests
    HAS_REQUESTS = True
except ImportError:
    HAS_REQUESTS = False

_MAX_CHARS = 200000
_OFFICE_EXTS = {
    ".pdf", ".docx", ".docm", ".xlsx", ".xlsm", ".xlsb", ".xltx", ".pptx", ".pptm", ".potx",
}

_TEXT_EXTS = {
    ".txt", ".text", ".md", ".markdown", ".rst", ".log", ".csv", ".tsv", ".json",
    ".jsonl", ".xml", ".html", ".htm", ".css", ".scss", ".less", ".js", ".jsx",
    ".ts", ".tsx", ".mjs", ".cjs", ".vue", ".svelte", ".py", ".pyw", ".pyi",
    ".java", ".kt", ".kts", ".c", ".h", ".cpp", ".hpp", ".cc", ".cs", ".go",
    ".rs", ".rb", ".php", ".sql", ".pl", ".pm", ".lua", ".r", ".sh", ".bash",
    ".bat", ".cmd", ".ps1", ".psm1", ".ini", ".cfg", ".conf", ".yaml", ".yml",
    ".toml", ".env", ".gitignore", ".dockerfile", ".properties", ".jsonc",
}

_CODE_EXTS = {e for e in _TEXT_EXTS if e not in {
    ".txt", ".text", ".md", ".markdown", ".rst", ".log", ".csv", ".tsv",
    ".json", ".jsonl", ".xml", ".html", ".htm", ".css", ".ini", ".cfg",
    ".conf", ".yaml", ".yml", ".toml", ".env", ".gitignore", ".properties",
    ".jsonc",
}}


def _detect_ext(path):
    p = str(path).lower()
    if p.endswith(".dockerfile") or Path(p).name.lower() == "dockerfile":
        return ".dockerfile"
    return Path(p).suffix.lower() if Path(p).suffix else ""


def _read_text(path, max_chars=_MAX_CHARS):
    """Lee un archivo de texto con deteccion de encoding."""
    data = Path(path).read_bytes()
    enc = None
    for codec in ("utf-8-sig", "utf-16", "utf-8", "cp1252", "latin-1"):
        try:
            data.decode(codec)
            enc = codec
            break
        except (UnicodeDecodeError, LookupError):
            continue
    if enc is None:
        enc = "latin-1"
    return data.decode(enc, errors="replace")[:max_chars]


def _read_pdf(path, max_chars=_MAX_CHARS):
    if not HAS_PDF:
        return None, "PyPDF2 no instalado"
    reader = PdfReader(str(path))
    parts = []
    total = 0
    for i, page in enumerate(reader.pages):
        txt = (page.extract_text() or "").strip()
        if txt:
            block = f"--- Pagina {i + 1} ---\n{txt}\n"
        else:
            block = f"--- Pagina {i + 1} (sin texto extraible) ---\n"
        parts.append(block)
        total += len(block)
        if total > max_chars:
            parts.append("\n[... texto truncado ...]")
            break
    return "\n".join(parts), f"{len(reader.pages)} paginas"


def _read_docx(path, max_chars=_MAX_CHARS):
    if not HAS_DOCX:
        return None, "python-docx no instalado"
    doc = _docx.Document(str(path))
    lines = []
    total = 0
    for para in doc.paragraphs:
        txt = para.text.rstrip()
        if txt:
            lines.append(txt)
            total += len(txt) + 1
            if total > max_chars:
                lines.append("[... texto truncado ...]")
                break
    for ti, table in enumerate(doc.tables[:10]):
        if total > max_chars:
            break
        lines.append(f"\n--- Tabla {ti + 1} ---")
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            lines.append(" | ".join(cells))
            total += len(" | ".join(cells)) + 1
    return "\n".join(lines), f"{len(doc.paragraphs)} parrafos, {len(doc.tables)} tablas"


def _read_xlsx(path, max_chars=_MAX_CHARS):
    if not HAS_XLSX:
        return None, "openpyxl no instalado"
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    lines = []
    total = 0
    for ws in wb.worksheets:
        if total > max_chars:
            break
        lines.append(f"\n=== Hoja: {ws.title} ===")
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            if not any(vals):
                continue
            line = " | ".join(vals)
            lines.append(line)
            total += len(line) + 1
            if total > max_chars:
                lines.append("[... texto truncado ...]")
                break
    wb.close()
    return "\n".join(lines), f"{len(wb.worksheets)} hojas"


def _read_pptx(path, max_chars=_MAX_CHARS):
    if not HAS_PPTX:
        return None, "python-pptx no instalado"
    prs = Presentation(str(path))
    lines = []
    total = 0
    for idx, slide in enumerate(prs.slides):
        if total > max_chars:
            break
        lines.append(f"\n=== Diapositiva {idx + 1} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        lines.append(txt)
                        total += len(txt) + 1
            elif shape.shape_type == 19 and shape.has_table:  # TABLE
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    line = " | ".join(cells)
                    lines.append(line)
                    total += len(line) + 1
            if total > max_chars:
                lines.append("[... texto truncado ...]")
                break
    return "\n".join(lines), f"{len(prs.slides)} diapositivas"


def _extract(path, max_chars=_MAX_CHARS):
    """Devuelve (texto, detalle_del_formato). None en texto si no soportado."""
    ext = _detect_ext(path)
    if ext in _TEXT_EXTS:
        try:
            return _read_text(path, max_chars), "texto"
        except Exception as e:
            return None, f"error leyendo texto: {e}"
    if ext == ".pdf":
        try:
            return _read_pdf(path, max_chars)
        except Exception as e:
            return None, f"error leyendo PDF: {e}"
    if ext in (".docx", ".docm"):
        try:
            return _read_docx(path, max_chars)
        except Exception as e:
            return None, f"error leyendo Word: {e}"
    if ext in (".xlsx", ".xlsm", ".xlsb", ".xltx"):
        try:
            return _read_xlsx(path, max_chars)
        except Exception as e:
            return None, f"error leyendo Excel: {e}"
    if ext in (".pptx", ".pptm", ".potx"):
        try:
            return _read_pptx(path, max_chars)
        except Exception as e:
            return None, f"error leyendo PowerPoint: {e}"
    return None, "formato no soportado"


def _is_text_file(path):
    return _detect_ext(path) in _TEXT_EXTS


def _format_name(ext):
    return {
        ".pdf": "PDF", ".docx": "Word (DOCX)", ".docm": "Word (DOCM)",
        ".xlsx": "Excel (XLSX)", ".xlsm": "Excel (XLSM)", ".xlsb": "Excel (XLSB)",
        ".xltx": "Excel (XLTX)", ".pptx": "PowerPoint (PPTX)",
        ".pptm": "PowerPoint (PPTM)", ".potx": "PowerPoint (POTX)",
    }.get(ext, "texto/codigo" if ext in _TEXT_EXTS else "desconocido")


def _format_payload(text):
    if len(text) <= 4000:
        return text
    return text[:4000] + "\n[...]"


def document_tool(parameters, player=None):
    """Dispatcher principal del tool de documentos."""
    action = str(parameters.get("action") or parameters.get("accion") or "read").lower()
    path = parameters.get("path") or parameters.get("file") or parameters.get("archivo") or ""
    if path:
        path = str(path).strip().strip('"').strip("'")

    if action in ("formats", "list", "formatos", "lista"):
        return (
            "Formatos soportados:\n"
            f"  Oficina: {', '.join(sorted(_OFFICE_EXTS))}\n"
            "  Texto/codigo: .txt .md .csv .json .xml .html .py .js .ts .java .c .cpp "
            ".go .rs .rb .php .sql .sh .bat .ps1 .yaml .toml .ini ...\n"
            "Acciones: info, read (extraer texto), summary, write, edit/replace, "
            "append, to_txt/convert."
        )

    if not path:
        return "Error: indica 'path' al archivo (y opcionalmente 'action')."
    if not os.path.isfile(path):
        return f"Error: archivo no encontrado: {path}"

    if action in ("info", "read-info", "metadatos"):
        st = os.stat(path)
        ext = _detect_ext(path)
        size = st.st_size
        mtime = datetime.fromtimestamp(st.st_mtime).strftime("%Y-%m-%d %H:%M")
        lines = [
            f"Archivo: {Path(path).name}",
            f"Ruta: {os.path.abspath(path)}",
            f"Tamano: {size:,} bytes ({size / 1024:.1f} KB)",
            f"Formato: {_format_name(ext)}",
            f"Ultima modificacion: {mtime}",
            f"Extensiones soportadas: {'SI' if ext in _OFFICE_EXTS or ext in _TEXT_EXTS else 'NO'}",
        ]
        return "\n".join(lines)

    if action in ("read", "extract", "leer", "contenido"):
        max_chars = int(parameters.get("max_chars", _MAX_CHARS) or _MAX_CHARS)
        text, detail = _extract(path, max_chars)
        if text is None:
            return f"Error: no se pudo leer el documento ({detail})."
        return (
            f"=== CONTENIDO DE '{Path(path).name}' ({detail}) ===\n"
            f"{_format_payload(text)}"
        )

    if action in ("summary", "resumen", "preview", "vista"):
        max_chars = int(parameters.get("max_chars", 12000) or 12000)
        text, detail = _extract(path, max_chars)
        if text is None:
            return f"Error: no se pudo leer el documento ({detail})."
        words = len(text.split())
        chars = len(text)
        return (
            f"=== RESUMEN/PREVIEW DE '{Path(path).name}' ===\n"
            f"Formato: {detail} | {chars:,} caracteres | ~{words:,} palabras\n\n"
            f"{_format_payload(text)}"
        )

    if action in ("write", "save", "guardar", "crear"):
        content = parameters.get("content") or parameters.get("text") or ""
        if not _is_text_file(path):
            return (
                f"Error: 'write' solo funciona con archivos de texto. '{Path(path).name}' "
                "es un formato binario de oficina. Usa 'to_txt' para convertirlo a texto."
            )
        if not content:
            return "Error: indica 'content' con el contenido a escribir."
        Path(path).write_text(content, encoding="utf-8")
        return f"Archivo guardado: {path} ({len(content):,} caracteres)."

    if action in ("edit", "replace", "reemplazar", "corregir"):
        find = parameters.get("find")
        replace = parameters.get("replace")
        if not _is_text_file(path):
            return (
                f"Error: 'edit' solo funciona con archivos de texto. "
                "Para Word/PDF/Excel/PPT extrae el contenido con 'read' y guarda "
                "la version corregida como texto con 'write' o 'to_txt'."
            )
        if not find:
            return "Error: indica 'find' con el texto a buscar."
        text = _read_text(path)
        if find not in text:
            return f"No se encontro la cadena: {find[:80]}"
        n = text.count(find)
        new_text = text.replace(find, replace or "")
        Path(path).write_text(new_text, encoding="utf-8")
        return (
            f"Reemplazadas {n} ocurrencia(s) de '{find[:60]}' por "
            f"'{str(replace)[:60]}' en {path}."
        )

    if action in ("append", "agregar"):
        content = parameters.get("content") or parameters.get("text") or ""
        if not _is_text_file(path):
            return f"Error: 'append' solo funciona con archivos de texto ({Path(path).name})."
        if not content:
            return "Error: indica 'content' con el texto a agregar."
        with open(path, "a", encoding="utf-8") as f:
            f.write(content if content.endswith("\n") else content + "\n")
        return f"Texto agregado a: {path}"

    if action in ("to_txt", "convert", "convertir", "exportar"):
        text, detail = _extract(path)
        if text is None:
            return f"Error: no se pudo convertir el documento ({detail})."
        out_path = str(Path(path).with_suffix(".txt"))
        if out_path.lower() == str(path).lower():
            out_path = str(path) + ".txt"
        Path(out_path).write_text(text, encoding="utf-8")
        return (
            f"Documento convertido de {_format_name(_detect_ext(path))} a texto plano.\n"
            f"Guardado en: {out_path} ({len(text):,} caracteres).\n"
            f"Contenido:\n{_format_payload(text)}"
        )

    return ("Tool de documentos disponible. Acciones:\n"
            "  info <path>  - metadatos del archivo\n"
            "  read <path> [max_chars]  - extrae el contenido (PDF/Word/Excel/PPT/texto/codigo)\n"
            "  summary <path> [max_chars]  - preview y conteo del contenido\n"
            "  write <path> content=<texto>  - escribe/sobrescribe archivo de texto\n"
            "  edit <path> find=<texto> [replace=<texto>]  - reemplaza en archivo de texto\n"
            "  append <path> content=<texto>  - agrega texto al final\n"
            "  to_txt <path>  - convierte un documento de oficina a .txt\n"
            "  formats  - lista los formatos soportados")
