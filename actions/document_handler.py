"""
document_handler.py — Manejo completo de documentos: Word, PowerPoint, Excel, PDF, Text.
Crear, leer, escribir, convertir, resumir, traducir, interpretar, abrir.
ERIS recuerda todo lo que creo.
"""
import json
import os
import subprocess
import time
from pathlib import Path
from datetime import datetime

_DATA_DIR = Path(__file__).resolve().parent.parent / "data"
_DOCS_DIR = _DATA_DIR / "documents"
_CREATED_FILE = _DATA_DIR / "documents_created.json"
_HISTORY_FILE = _DATA_DIR / "document_history.json"

_FORMATS = {
    ".docx": "Word",
    ".pptx": "PowerPoint",
    ".xlsx": "Excel",
    ".pdf": "PDF",
    ".txt": "Texto",
    ".csv": "CSV",
    ".md": "Markdown",
    ".html": "HTML",
    ".json": "JSON",
    ".py": "Python",
    ".js": "JavaScript",
}


def document_handler(parameters: dict = None, player=None) -> str:
    """Manejo completo de documentos con memoria."""
    params = parameters or {}
    action = params.get("action", "help").lower()

    if action == "create_word":
        return _create_word(params)
    elif action == "create_pptx":
        return _create_pptx(params)
    elif action == "create_excel":
        return _create_excel(params)
    elif action == "create_pdf":
        return _create_pdf(params)
    elif action == "create_txt":
        return _create_txt(params)
    elif action == "create_csv":
        return _create_csv(params)
    elif action == "read":
        return _read_any(params)
    elif action == "read_word":
        return _read_word(params)
    elif action == "read_pptx":
        return _read_pptx(params)
    elif action == "read_excel":
        return _read_excel(params)
    elif action == "read_pdf":
        return _read_pdf(params)
    elif action == "read_txt":
        return _read_txt(params)
    elif action == "convert_to_pdf":
        return _convert_to_pdf(params)
    elif action == "merge_pdfs":
        return _merge_pdfs(params)
    elif action == "split_pdf":
        return _split_pdf(params)
    elif action == "info":
        return _get_info(params)
    elif action == "open":
        return _open_doc(params)
    elif action == "summarize":
        return _summarize(params)
    elif action == "translate":
        return _translate(params)
    elif action == "interpret":
        return _interpret(params)
    elif action == "what_i_wrote":
        return _what_i_wrote()
    elif action == "working_doc":
        return _get_working_doc()
    elif action == "list_recent":
        return _list_recent()
    elif action == "clear_history":
        return _clear_history()
    return (
        "CREAR: create_word, create_pptx, create_excel, create_pdf, create_txt, create_csv\n"
        "LEER: read (auto), read_word, read_pptx, read_excel, read_pdf, read_txt\n"
        "CONVERTIR: convert_to_pdf, merge_pdfs, split_pdf\n"
        "ANALIZAR: summarize, translate, interpret, info\n"
        "OTROS: open, what_i_wrote, working_doc, list_recent, clear_history"
    )


# ═══════════════════════════════════════════════════════════════
#  CREAR DOCUMENTOS
# ═══════════════════════════════════════════════════════════════

def _create_word(params: dict) -> str:
    title = params.get("title", "Documento")
    content = params.get("content", "")
    save_path = _resolve_path(params.get("path", ""), ".docx")
    try:
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        doc = Document()
        style = doc.styles["Normal"]
        style.font.size = Pt(11)
        style.font.name = "Calibri"
        doc.add_heading(title, 0)
        for line in content.split("\n"):
            line = line.strip()
            if not line:
                doc.add_paragraph("")
            elif line.startswith("# "):
                doc.add_heading(line[2:], level=1)
            elif line.startswith("## "):
                doc.add_heading(line[3:], level=2)
            elif line.startswith("### "):
                doc.add_heading(line[4:], level=3)
            elif line.startswith("- ") or line.startswith("* "):
                doc.add_paragraph(line[2:], style="List Bullet")
            elif line.startswith("1. ") or line.startswith("2. "):
                doc.add_paragraph(line[3:], style="List Number")
            elif line.startswith("**") and line.endswith("**"):
                p = doc.add_paragraph()
                run = p.add_run(line[2:-2])
                run.bold = True
            else:
                doc.add_paragraph(line)
        doc.save(str(save_path))
        _track_created("Word", save_path, title)
        return "Word creado: '{}'".format(save_path)
    except Exception as e:
        return "Error creando Word: {}".format(str(e)[:80])


def _create_pptx(params: dict) -> str:
    title = params.get("title", "Presentacion")
    content = params.get("content", "")
    save_path = _resolve_path(params.get("path", ""), ".pptx")
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content[:500] if content else "Contenido"
        sections = content.split("\n---\n")
        for section in sections[1:6]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            lines = section.strip().split("\n")
            if lines:
                slide.shapes.title.text = lines[0][:60]
                slide.placeholders[1].text = "\n".join(lines[1:])[:500]
        prs.save(str(save_path))
        _track_created("PowerPoint", save_path, title)
        return "PowerPoint creado: '{}' ({} slides)".format(save_path, len(prs.slides))
    except Exception as e:
        return "Error creando PowerPoint: {}".format(str(e)[:80])


def _create_excel(params: dict) -> str:
    title = params.get("title", "Hoja de calculo")
    content = params.get("content", "")
    save_path = _resolve_path(params.get("path", ""), ".xlsx")
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = title[:31]
        header_font = Font(bold=True, size=12, color="FFFFFF")
        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        if content:
            lines = content.strip().split("\n")
            for row_idx, line in enumerate(lines, 1):
                cells = [c.strip() for c in line.split(",")]
                for col_idx, cell in enumerate(cells, 1):
                    ws.cell(row=row_idx, column=col_idx, value=cell)
                    if row_idx == 1:
                        cell_obj = ws.cell(row=row_idx, column=col_idx)
                        cell_obj.font = header_font
                        cell_obj.fill = header_fill
        else:
            ws.cell(row=1, column=1, value="Columna A")
            ws.cell(row=1, column=2, value="Columna B")
            ws.cell(row=1, column=3, value="Columna C")
        for col in ws.columns:
            max_length = 0
            col_letter = col[0].column_letter
            for cell in col:
                if cell.value:
                    max_length = max(max_length, len(str(cell.value)))
            ws.column_dimensions[col_letter].width = min(max_length + 4, 50)
        wb.save(str(save_path))
        _track_created("Excel", save_path, title)
        return "Excel creado: '{}' ({} filas)".format(save_path, ws.max_row)
    except Exception as e:
        return "Error creando Excel: {}".format(str(e)[:80])


def _create_pdf(params: dict) -> str:
    title = params.get("title", "Documento")
    content = params.get("content", "")
    save_path = _resolve_path(params.get("path", ""), ".pdf")
    try:
        from fpdf import FPDF
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 20)
        pdf.multi_cell(w=0, h=15, text=title)
        pdf.set_x(10)
        pdf.ln(3)
        pdf.set_font("Helvetica", size=11)
        for line in content.split("\n"):
            line = line.strip()
            if not line:
                pdf.ln(3)
                pdf.set_x(10)
            elif line.startswith("# "):
                pdf.set_font("Helvetica", "B", 16)
                pdf.multi_cell(w=0, h=10, text=line[2:])
                pdf.set_x(10)
                pdf.set_font("Helvetica", size=11)
            elif line.startswith("## "):
                pdf.set_font("Helvetica", "B", 13)
                pdf.multi_cell(w=0, h=8, text=line[3:])
                pdf.set_x(10)
                pdf.set_font("Helvetica", size=11)
            elif line.startswith("- "):
                pdf.multi_cell(w=0, h=7, text="  * " + line[2:])
                pdf.set_x(10)
            else:
                pdf.multi_cell(w=0, h=7, text=line)
                pdf.set_x(10)
        pdf.output(str(save_path))
        _track_created("PDF", save_path, title)
        return "PDF creado: '{}'".format(save_path)
    except Exception as e:
        return "Error creando PDF: {}".format(str(e)[:80])


def _create_txt(params: dict) -> str:
    title = params.get("title", "documento")
    content = params.get("content", "")
    save_path = _resolve_path(params.get("path", ""), ".txt")
    try:
        save_path.parent.mkdir(parents=True, exist_ok=True)
        save_path.write_text(content, encoding="utf-8")
        _track_created("Texto", save_path, title)
        return "Texto creado: '{}' ({} caracteres)".format(save_path, len(content))
    except Exception as e:
        return "Error creando texto: {}".format(str(e)[:80])


def _create_csv(params: dict) -> str:
    content = params.get("content", "")
    save_path = _resolve_path(params.get("path", ""), ".csv")
    try:
        save_path.parent.mkdir(parents=True, exist_ok=True)
        save_path.write_text(content, encoding="utf-8")
        _track_created("CSV", save_path, "CSV")
        return "CSV creado: '{}'".format(save_path)
    except Exception as e:
        return "Error creando CSV: {}".format(str(e)[:80])


# ═══════════════════════════════════════════════════════════════
#  LEER DOCUMENTOS
# ═══════════════════════════════════════════════════════════════

def _read_any(params: dict) -> str:
    path = params.get("path", "")
    if not path:
        return "Error: especifica 'path' del archivo"
    p = Path(path)
    if not p.exists():
        return "Archivo no encontrado: '{}'".format(path)
    ext = p.suffix.lower()
    if ext == ".docx":
        return _read_word(params)
    elif ext == ".pptx":
        return _read_pptx(params)
    elif ext == ".xlsx" or ext == ".xls":
        return _read_excel(params)
    elif ext == ".pdf":
        return _read_pdf(params)
    elif ext == ".txt" or ext == ".md" or ext == ".log":
        return _read_txt(params)
    elif ext == ".csv":
        return _read_txt(params)
    elif ext == ".json":
        return _read_txt(params)
    elif ext == ".html" or ext == ".htm":
        return _read_txt(params)
    elif ext == ".py" or ext == ".js" or ext == ".ts":
        return _read_txt(params)
    else:
        return _read_txt(params)


def _read_word(params: dict) -> str:
    path = params.get("path", "")
    if not path:
        return "Error: especifica 'path'"
    try:
        from docx import Document
        doc = Document(path)
        lines = []
        for para in doc.paragraphs:
            style = para.style.name
            text = para.text.strip()
            if not text:
                continue
            if "Heading 1" in style:
                lines.append("# " + text)
            elif "Heading 2" in style:
                lines.append("## " + text)
            elif "Heading 3" in style:
                lines.append("### " + text)
            elif "List Bullet" in style:
                lines.append("- " + text)
            elif "List Number" in style:
                lines.append("* " + text)
            else:
                lines.append(text)
        tables = []
        for i, table in enumerate(doc.tables):
            tables.append("\n--- TABLA {} ---".format(i + 1))
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                tables.append(" | ".join(cells))
        result = "═══ WORD: {} ═══\n\n".format(Path(path).name)
        result += "\n".join(lines)
        if tables:
            result += "\n\n" + "\n".join(tables)
        result += "\n\n---\n{} parrafos, {} tablas".format(len(doc.paragraphs), len(doc.tables))
        return result[:8000]
    except Exception as e:
        return "Error leyendo Word: {}".format(str(e)[:80])


def _read_pptx(params: dict) -> str:
    path = params.get("path", "")
    if not path:
        return "Error: especifica 'path'"
    try:
        from pptx import Presentation
        prs = Presentation(path)
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append("--- SLIDE {} ---".format(i))
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
            lines.append("")
        result = "═══ POWERPOINT: {} ═══\n\n".format(Path(path).name)
        result += "\n".join(lines)
        result += "\n\n---\n{} slides".format(len(prs.slides))
        return result[:8000]
    except Exception as e:
        return "Error leyendo PowerPoint: {}".format(str(e)[:80])


def _read_excel(params: dict) -> str:
    path = params.get("path", "")
    if not path:
        return "Error: especifica 'path'"
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append("--- HOJA: {} ---".format(sheet_name))
            row_count = 0
            for row in ws.iter_rows(max_row=50, values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                non_empty = [c for c in cells if c]
                if non_empty:
                    lines.append(" | ".join(cells))
                    row_count += 1
            lines.append("({} filas)".format(row_count))
            lines.append("")
        wb.close()
        result = "═══ EXCEL: {} ═══\n\n".format(Path(path).name)
        result += "\n".join(lines)
        result += "\n\n---\n{} hojas".format(len(wb.sheetnames))
        return result[:8000]
    except Exception as e:
        return "Error leyendo Excel: {}".format(str(e)[:80])


def _read_pdf(params: dict) -> str:
    path = params.get("path", "")
    if not path:
        return "Error: especifica 'path'"
    try:
        import fitz
        doc = fitz.open(path)
        lines = []
        for i, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                lines.append("--- PAGINA {} ---".format(i + 1))
                lines.append(text.strip())
        doc.close()
        result = "═══ PDF: {} ═══\n\n".format(Path(path).name)
        result += "\n".join(lines)
        result += "\n\n---\n{} paginas".format(len(doc))
        return result[:8000]
    except Exception as e:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(path)
            lines = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    lines.append("--- PAGINA {} ---".format(i + 1))
                    lines.append(text.strip())
            result = "═══ PDF: {} ═══\n\n".format(Path(path).name)
            result += "\n".join(lines)
            result += "\n\n---\n{} paginas".format(len(reader.pages))
            return result[:8000]
        except Exception as e2:
            return "Error leyendo PDF: {}".format(str(e2)[:80])


def _read_txt(params: dict) -> str:
    path = params.get("path", "")
    if not path:
        return "Error: especifica 'path'"
    try:
        p = Path(path)
        content = p.read_text(encoding="utf-8", errors="replace")
        result = "═══ {}: {} ═══\n\n".format(_FORMATS.get(p.suffix.lower(), "Archivo"), p.name)
        result += content
        result += "\n\n---\n{} caracteres, {} lineas".format(len(content), len(content.split("\n")))
        return result[:8000]
    except Exception as e:
        return "Error leyendo archivo: {}".format(str(e)[:80])


# ═══════════════════════════════════════════════════════════════
#  CONVERTIR
# ═══════════════════════════════════════════════════════════════

def _convert_to_pdf(params: dict) -> str:
    path = params.get("path", "")
    if not path:
        return "Error: especifica 'path'"
    p = Path(path)
    if not p.exists():
        return "Archivo no encontrado: '{}'".format(path)
    ext = p.suffix.lower()
    save_path = params.get("output", str(p.with_suffix(".pdf")))
    if ext == ".pdf":
        return "Ya es PDF"
    if ext == ".txt" or ext == ".md" or ext == ".csv" or ext == ".json" or ext == ".py":
        return _convert_text_to_pdf(p, save_path)
    elif ext == ".docx":
        return _convert_docx_to_pdf(p, save_path)
    elif ext == ".pptx":
        return _convert_pptx_to_pdf(p, save_path)
    elif ext == ".xlsx":
        return _convert_xlsx_to_pdf(p, save_path)
    else:
        return "Formato no soportado para conversion: {}".format(ext)


def _convert_text_to_pdf(src: Path, save_path: str) -> str:
    try:
        from fpdf import FPDF
        content = src.read_text(encoding="utf-8", errors="replace")
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Helvetica", "", 11)
        for line in content.split("\n"):
            pdf.multi_cell(w=0, h=7, text=line)
            pdf.set_x(10)
        pdf.output(save_path)
        _track_created("PDF (convertido)", Path(save_path), src.name)
        return "Convertido a PDF: '{}'".format(save_path)
    except Exception as e:
        return "Error: {}".format(str(e)[:60])


def _convert_docx_to_pdf(src: Path, save_path: str) -> str:
    try:
        from docx import Document
        from fpdf import FPDF
        doc = Document(str(src))
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Helvetica", "", 11)
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                pdf.ln(3)
                continue
            style = para.style.name
            if "Heading 1" in style:
                pdf.set_font("Helvetica", "B", 16)
                pdf.multi_cell(w=0, h=10, text=text)
                pdf.set_x(10)
                pdf.set_font("Helvetica", "", 11)
            elif "Heading 2" in style:
                pdf.set_font("Helvetica", "B", 13)
                pdf.multi_cell(w=0, h=8, text=text)
                pdf.set_x(10)
                pdf.set_font("Helvetica", "", 11)
            else:
                pdf.multi_cell(w=0, h=7, text=text)
                pdf.set_x(10)
        pdf.output(save_path)
        _track_created("PDF (desde Word)", Path(save_path), src.name)
        return "Word convertido a PDF: '{}'".format(save_path)
    except Exception as e:
        return "Error: {}".format(str(e)[:60])


def _convert_pptx_to_pdf(src: Path, save_path: str) -> str:
    try:
        from pptx import Presentation
        from fpdf import FPDF
        prs = Presentation(str(src))
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        for i, slide in enumerate(prs.slides, 1):
            pdf.add_page()
            pdf.set_font("Helvetica", "B", 14)
            pdf.multi_cell(w=0, h=10, text="Slide {}".format(i))
            pdf.set_x(10)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            pdf.set_font("Helvetica", "", 11)
                            pdf.multi_cell(w=0, h=7, text=text)
                            pdf.set_x(10)
        pdf.output(save_path)
        _track_created("PDF (desde PowerPoint)", Path(save_path), src.name)
        return "PowerPoint convertido a PDF: '{}'".format(save_path)
    except Exception as e:
        return "Error: {}".format(str(e)[:60])


def _convert_xlsx_to_pdf(src: Path, save_path: str) -> str:
    try:
        import openpyxl
        from fpdf import FPDF
        wb = openpyxl.load_workbook(str(src), read_only=True, data_only=True)
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            pdf.add_page()
            pdf.set_font("Helvetica", "B", 12)
            pdf.multi_cell(w=0, h=10, text="Hoja: {}".format(sheet_name))
            pdf.set_x(10)
            pdf.set_font("Helvetica", "", 9)
            for row in ws.iter_rows(max_row=100, values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = " | ".join(cells)
                pdf.multi_cell(w=0, h=6, text=line[:200])
                pdf.set_x(10)
        wb.close()
        pdf.output(save_path)
        _track_created("PDF (desde Excel)", Path(save_path), src.name)
        return "Excel convertido a PDF: '{}'".format(save_path)
    except Exception as e:
        return "Error: {}".format(str(e)[:60])


def _merge_pdfs(params: dict) -> str:
    paths = params.get("paths", [])
    save_path = params.get("output", str(_DOCS_DIR / "merged.pdf"))
    if not paths or len(paths) < 2:
        return "Error: especifica 'paths' con al menos 2 PDFs"
    try:
        from PyPDF2 import PdfMerger
        merger = PdfMerger()
        for p in paths:
            if Path(p).exists():
                merger.append(p)
        merger.write(save_path)
        merger.close()
        _track_created("PDF (fusionado)", Path(save_path), "{} PDFs".format(len(paths)))
        return "PDFs fusionados: '{}' ({} archivos)".format(save_path, len(paths))
    except Exception as e:
        return "Error fusionando: {}".format(str(e)[:60])


def _split_pdf(params: dict) -> str:
    path = params.get("path", "")
    if not path:
        return "Error: especifica 'path'"
    try:
        from PyPDF2 import PdfReader, PdfWriter
        reader = PdfReader(path)
        out_dir = Path(params.get("output_dir", str(_DOCS_DIR / "split")))
        out_dir.mkdir(parents=True, exist_ok=True)
        count = 0
        for i, page in enumerate(reader.pages):
            writer = PdfWriter()
            writer.add_page(page)
            out_path = out_dir / "page_{}.pdf".format(i + 1)
            writer.write(str(out_path))
            count += 1
        return "PDF dividido: {} paginas en '{}'".format(count, out_dir)
    except Exception as e:
        return "Error: {}".format(str(e)[:60])


# ═══════════════════════════════════════════════════════════════
#  INFO Y ABRIR
# ═══════════════════════════════════════════════════════════════

def _get_info(params: dict) -> str:
    path = params.get("path", "")
    if not path:
        created = _load_created()
        files = list(_DOCS_DIR.glob("*")) if _DOCS_DIR.exists() else []
        total_kb = sum(f.stat().st_size / 1024 for f in files if f.is_file())
        lines = [
            "═══ SISTEMA DE DOCUMENTOS ═══",
            "",
            "  Documentos creados: {}".format(len(created)),
            "  Archivos en carpeta: {}".format(len(files)),
            "  Tamano total: {:.1f} KB".format(total_kb),
            "",
            "  Carpeta: {}".format(str(_DOCS_DIR)),
            "",
            "  Formatos soportados:",
            "    CREAR: Word (.docx), Excel (.xlsx), PowerPoint (.pptx), PDF, TXT, CSV",
            "    LEER:  Word, Excel, PowerPoint, PDF, TXT, CSV, JSON",
            "    CONVERTIR: TXT/DOCX/PPTX/XLSX -> PDF",
            "    ANALIZAR: Resumir, traducir (stub), interpretar (stub)",
            "    MEMORIA: Que escribi, listar recientes, limpiar historial",
        ]
        return "\n".join(lines)
    p = Path(path)
    if not p.exists():
        return "No encontrado: '{}'".format(path)
    stat = p.stat()
    ext = p.suffix.lower()
    fmt = _FORMATS.get(ext, "Desconocido")
    size_kb = stat.st_size / 1024
    mod = datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M")
    lines = [
        "═══ INFO DEL ARCHIVO ═══",
        "",
        "  Nombre:   {}".format(p.name),
        "  Formato:  {}".format(fmt),
        "  Tamano:   {:.1f} KB".format(size_kb),
        "  Modificado: {}".format(mod),
        "  Ruta:     {}".format(str(p)),
    ]
    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(str(p))
            lines.append("  Parrafos: {}".format(len(doc.paragraphs)))
            lines.append("  Tablas:   {}".format(len(doc.tables)))
        except Exception:
            pass
    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(p))
            lines.append("  Slides:   {}".format(len(prs.slides)))
        except Exception:
            pass
    elif ext == ".xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(p), read_only=True)
            lines.append("  Hojas:    {}".format(len(wb.sheetnames)))
            wb.close()
        except Exception:
            pass
    elif ext == ".pdf":
        try:
            import fitz
            doc = fitz.open(str(p))
            lines.append("  Paginas:  {}".format(len(doc)))
            meta = doc.metadata
            if meta:
                if meta.get("title"):
                    lines.append("  Titulo:   {}".format(meta["title"]))
                if meta.get("author"):
                    lines.append("  Autor:    {}".format(meta["author"]))
            doc.close()
        except Exception:
            pass
    return "\n".join(lines)


def _open_doc(params: dict) -> str:
    path = params.get("path", "")
    if not path:
        return "Error: especifica 'path'"
    p = Path(path)
    if not p.exists():
        return "No encontrado: '{}'".format(path)
    try:
        os.startfile(str(p))
        _track_opened(p)
        return "Abriendo: '{}'".format(p.name)
    except Exception as e:
        return "Error abriendo: {}".format(str(e)[:50])


# ═══════════════════════════════════════════════════════════════
#  ANALIZAR: RESUMIR, TRADUCIR, INTERPRETAR
# ═══════════════════════════════════════════════════════════════

def _summarize(params: dict) -> str:
    path = params.get("path", "")
    content = params.get("content", "")
    if path:
        read_params = {"path": path}
        raw = _read_any(read_params)
        content = _extract_text_from_result(raw)
    if not content:
        return "Error: especifica 'path' o 'content'"
    sentences = [s.strip() for s in content.replace("\n", " ").split(".") if len(s.strip()) > 20]
    if not sentences:
        return "Contenido muy corto para resumir"
    important = sentences[:min(10, len(sentences))]
    summary = ". ".join(important) + "."
    word_count = len(summary.split())
    return (
        "═══ RESUMEN ({}) ═══\n\n{}\n\n---\n{} palabras de {} originales".format(
            Path(path).name if path else "texto", summary[:2000], word_count, len(content.split())
        )
    )


def _translate(params: dict) -> str:
    content = params.get("content", "")
    target_lang = params.get("target_lang", "es")
    path = params.get("path", "")
    if path:
        read_params = {"path": path}
        raw = _read_any(read_params)
        content = _extract_text_from_result(raw)
    if not content:
        return "Error: especifica 'path' o 'content'"
    lang_names = {"es": "espanol", "en": "ingles", "fr": "frances", "de": "aleman",
                  "pt": "portugues", "it": "italiano", "ja": "japones", "zh": "chino",
                  "ko": "coreano", "ru": "ruso", "ar": "arabe"}
    target_name = lang_names.get(target_lang, target_lang)
    lines = [
        "═══ TRADUCCION -> {} ═══".format(target_name.upper()),
        "",
        "Para traducir contenido completo, necesito",
        "usar un modelo de IA local (Ollama) o remoto.",
        "",
        "Contenido original ({} caracteres):".format(len(content)),
        content[:1500],
        "",
        "---",
        "Traduccion pendiente — usar con modelo IA"
    ]
    return "\n".join(lines)


def _interpret(params: dict) -> str:
    path = params.get("path", "")
    content = params.get("content", "")
    question = params.get("question", "Que contiene este documento?")
    if path:
        read_params = {"path": path}
        raw = _read_any(read_params)
        content = _extract_text_from_result(raw)
    if not content:
        return "Error: especifica 'path' o 'content'"
    word_count = len(content.split())
    lines = [
        "═══ INTERPRETACION ═══",
        "",
        "Pregunta: {}".format(question),
        "",
        "Documento: {} ({} palabras)".format(Path(path).name if path else "texto", word_count),
        "",
        "Contenido detectado:",
        "- Formato: {}".format(Path(path).suffix if path else "texto plano"),
        "- Tamano: {} caracteres".format(len(content)),
        "- Lineas: {}".format(len(content.split("\n"))),
        "",
        "Primeras lineas:",
        content[:500],
        "",
        "---",
        "Para interpretacion profunda, usar con modelo IA (Ollama/local)"
    ]
    return "\n".join(lines)


# ═══════════════════════════════════════════════════════════════
#  MEMORIA: QUE CREO ERIS
# ═══════════════════════════════════════════════════════════════

def _what_i_wrote() -> str:
    created = _load_created()
    if not created:
        return "No he creado ningun documento"
    lines = ["═══ DOCUMENTOS QUE YO CREE ({}) ═══".format(len(created)), ""]
    for c in created[-20:]:
        lines.append("  [{}] {} '{}' ({})".format(
            c.get("time", "?")[:16],
            c.get("format", "?"),
            c.get("title", "?")[:40],
            str(Path(c.get("path", "?")).name)[:30]))
    return "\n".join(lines)


def _list_recent() -> str:
    if not _DOCS_DIR.exists():
        return "No hay documentos"
    files = sorted(_DOCS_DIR.glob("*"), key=lambda f: f.stat().st_mtime, reverse=True)
    if not files:
        return "No hay documentos en data/documents/"
    lines = ["═══ DOCUMENTOS RECIENTES ═══", ""]
    for f in files[:15]:
        fmt = _FORMATS.get(f.suffix.lower(), "?")
        size = f.stat().st_size / 1024
        mod = datetime.fromtimestamp(f.stat().st_mtime).strftime("%m-%d %H:%M")
        lines.append("  [{}] {:12s} {:8.1f} KB  {}".format(mod, fmt, size, f.name))
    return "\n".join(lines)


def _clear_history() -> str:
    _save_created([])
    _save_history([])
    return "Historial de documentos limpiado"


def _get_working_doc() -> str:
    _base = Path(__file__).resolve().parent.parent / "data"
    _wf = _base / "working_document.json"
    if not _wf.exists():
        return "No hay documento en curso."
    try:
        import json as _json
        data = _json.loads(_wf.read_text(encoding="utf-8"))
        return "Documento en curso:\n  Titulo: {}\n  Ruta: {}\n  Ultima modificacion: {}".format(
            data.get("title", "?"), data.get("path", "?"), data.get("time", "?")[:16])
    except Exception:
        return "Error leyendo documento en curso."


# ═══════════════════════════════════════════════════════════════
#  UTILIDADES
# ═══════════════════════════════════════════════════════════════

def _resolve_path(path_str: str, default_ext: str) -> Path:
    if path_str:
        p = Path(path_str)
        # If path is a directory (no file extension), auto-generate filename inside it
        if p.suffix not in (".docx", ".pptx", ".xlsx", ".pdf", ".txt", ".csv", ".md", ".html", ".json", ".py", ".js"):
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            p = p / "doc_{}{}".format(ts, default_ext)
    else:
        # Default: save to Desktop/ERIS_Documentos
        try:
            from actions.path_helper import get_desktop_path
            desk = Path(get_desktop_path())
        except Exception:
            desk = Path.home() / "Desktop"
        p = desk / "ERIS_Documentos" / "doc_{}{}".format(datetime.now().strftime("%Y%m%d_%H%M%S"), default_ext)
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def _extract_text_from_result(result: str) -> str:
    lines = result.split("\n")
    content_lines = []
    in_content = False
    for line in lines:
        if line.startswith("═══"):
            in_content = True
            continue
        if line.startswith("---") and in_content:
            break
        if in_content:
            content_lines.append(line)
    return "\n".join(content_lines).strip()


def _track_created(fmt: str, path: Path, title: str):
    created = _load_created()
    created.append({
        "format": fmt,
        "path": str(path),
        "title": title[:80],
        "time": datetime.now().isoformat(),
    })
    if len(created) > 100:
        created = created[-100:]
    _save_created(created)
    _log_history("create", str(path), fmt)
    try:
        import json as _json
        _base = Path(__file__).resolve().parent.parent / "data"
        _wf = _base / "working_document.json"
        _wf.parent.mkdir(parents=True, exist_ok=True)
        _wf.write_text(_json.dumps({"title": title, "path": str(path), "time": datetime.now().isoformat()}, indent=2, ensure_ascii=False), encoding="utf-8")
    except Exception:
        pass


def _track_opened(path: Path):
    _log_history("open", str(path), _FORMATS.get(path.suffix.lower(), "?"))


def _log_history(action: str, target: str, fmt: str):
    history = _load_history()
    history.append({
        "action": action,
        "target": target[:100],
        "format": fmt,
        "time": datetime.now().isoformat(),
    })
    if len(history) > 200:
        history = history[-200:]
    _save_history(history)


def _load_created() -> list:
    if _CREATED_FILE.exists():
        try:
            return json.loads(_CREATED_FILE.read_text(encoding="utf-8"))
        except Exception:
            pass
    return []


def _save_created(data: list):
    _DATA_DIR.mkdir(parents=True, exist_ok=True)
    _CREATED_FILE.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")


def _load_history() -> list:
    if _HISTORY_FILE.exists():
        try:
            return json.loads(_HISTORY_FILE.read_text(encoding="utf-8"))
        except Exception:
            pass
    return []


def _save_history(data: list):
    _DATA_DIR.mkdir(parents=True, exist_ok=True)
    _HISTORY_FILE.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
