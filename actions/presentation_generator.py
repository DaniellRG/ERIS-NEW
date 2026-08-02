# -*- coding: utf-8 -*-
"""presentation_generator.py — Crea presentaciones PowerPoint profesionales.
Soporta: portada, diapositivas con viñetas, lineas de tiempo, citas,
fondos de color, titulos, subtitulos. Crea carpeta automaticamente."""
import json
import os
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False


def presentation_generator(parameters: dict, player=None) -> str:
    action = parameters.get("action", "create").lower()
    if action == "create":
        return _create_presentation(parameters, player)
    elif action == "list_templates":
        return _list_templates()
    return "Acciones: create (crear pptx), list_templates (ayuda)."


def _create_presentation(parameters: dict, player=None) -> str:
    title = parameters.get("title", "Presentacion")
    subtitle = parameters.get("subtitle", "")
    author = parameters.get("author", "ERIS")
    slides_raw = parameters.get("slides", "[]")
    filename = parameters.get("filename", "")
    output_path_raw = parameters.get("output_path", "")

    if not _PPTX_OK:
        return "Error: Falta python-pptx. Instala con: pip install python-pptx"

    try:
        slides = json.loads(slides_raw) if isinstance(slides_raw, str) else slides_raw
    except json.JSONDecodeError:
        return "Error: slides debe ser un JSON valido."

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]

        _add_title_slide(prs, blank_layout, title, subtitle, author)

        for s in slides:
            slide_type = s.get("type", "content")
            slide_title = s.get("title", "")
            items = s.get("items", [])
            bg = s.get("bg_color", "")
            content = s.get("content", "")

            if slide_type == "section":
                _add_section_slide(prs, blank_layout, slide_title, bg)
            elif slide_type == "bullets":
                _add_bullet_slide(prs, blank_layout, slide_title, items, bg)
            elif slide_type == "numbered":
                _add_numbered_slide(prs, blank_layout, slide_title, items, bg)
            elif slide_type == "timeline":
                _add_timeline_slide(prs, blank_layout, slide_title, items, bg)
            elif slide_type == "quote":
                _add_quote_slide(prs, blank_layout, slide_title, content, bg)
            elif slide_type == "thankyou":
                _add_thankyou_slide(prs, blank_layout, slide_title, content)
            else:
                _add_bullet_slide(prs, blank_layout, slide_title, items, bg)

        output_dir = _resolve_output_dir(output_path_raw)
        output_dir.mkdir(parents=True, exist_ok=True)

        safe_name = "".join(c for c in title if c.isalnum() or c in " _-").strip().replace(" ", "_")
        if not filename:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{safe_name}_{ts}.pptx"
        if not filename.endswith(".pptx"):
            filename += ".pptx"

        output_path = output_dir / filename
        prs.save(str(output_path))

        if player:
            player.write_log(f"Presentacion creada: {output_path} ({len(prs.slides)} diapositivas)")

        return f"Presentacion creada: {output_path} ({len(prs.slides)} diapositivas, {os.path.getsize(output_path)} bytes)"

    except Exception as e:
        return f"Error creando presentacion: {e}"


def _add_bg(slide, color_hex: str):
    if not color_hex:
        return
    try:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        r = int(color_hex[0:2], 16)
        g = int(color_hex[2:4], 16)
        b = int(color_hex[4:6], 16)
        fill.fore_color.rgb = RGBColor(r, g, b)
    except Exception:
        pass


def _add_text_box(slide, left, top, width, height, text, size=18, bold=False, color=None, align=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = 'Calibri'
    if color:
        p.font.color.rgb = color
    if align:
        from pptx.enum.text import PP_ALIGN
        p.alignment = align
    return txBox


def _add_title_slide(prs, layout, title, subtitle, author):
    slide = prs.slides.add_slide(layout)
    _add_bg(slide, "1A3C6E")
    from pptx.enum.text import PP_ALIGN
    _add_text_box(slide, 1, 2, 11, 2, title, 44, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.CENTER)
    if subtitle:
        _add_text_box(slide, 1, 4.5, 11, 1, subtitle, 24, False, RGBColor(0xAA, 0xCC, 0xEE), PP_ALIGN.CENTER)
    _add_text_box(slide, 1, 6.5, 11, 0.5, f"Por: {author}", 14, False, RGBColor(0x88, 0xAA, 0xCC), PP_ALIGN.CENTER)


def _add_section_slide(prs, layout, title, bg):
    slide = prs.slides.add_slide(layout)
    _add_bg(slide, bg or "1A3C6E")
    from pptx.enum.text import PP_ALIGN
    _add_text_box(slide, 1, 2.5, 11, 2, title, 40, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.CENTER)


def _add_bullet_slide(prs, layout, title, items, bg):
    slide = prs.slides.add_slide(layout)
    _add_bg(slide, bg or "FFFFFF")
    _add_text_box(slide, 0.5, 0.3, 12, 1, title, 32, True, RGBColor(0x1A, 0x3C, 0x6E))
    for i, item in enumerate(items):
        y = 1.8 + i * 0.85
        _add_text_box(slide, 0.8, y, 11, 0.7, f"  {item}", 16, False, RGBColor(0x33, 0x33, 0x33))


def _add_numbered_slide(prs, layout, title, items, bg):
    slide = prs.slides.add_slide(layout)
    _add_bg(slide, bg or "FFFFFF")
    _add_text_box(slide, 0.5, 0.3, 12, 1, title, 32, True, RGBColor(0x1A, 0x3C, 0x6E))
    for i, item in enumerate(items):
        y = 1.8 + i * 0.85
        _add_text_box(slide, 0.8, y, 11, 0.7, f"{i+1}. {item}", 16, False, RGBColor(0x33, 0x33, 0x33))


def _add_timeline_slide(prs, layout, title, items, bg):
    slide = prs.slides.add_slide(layout)
    _add_bg(slide, bg or "F0F0F0")
    _add_text_box(slide, 0.5, 0.3, 12, 1, title, 32, True, RGBColor(0x1A, 0x3C, 0x6E))
    for i, item in enumerate(items):
        y = 1.8 + i * 0.85
        if isinstance(item, dict):
            year = item.get("year", "")
            desc = item.get("desc", "")
            _add_text_box(slide, 0.8, y, 1.5, 0.6, year, 18, True, RGBColor(0x1A, 0x3C, 0x6E))
            _add_text_box(slide, 2.5, y, 9, 0.6, desc, 14, False, RGBColor(0x33, 0x33, 0x33))
        else:
            _add_text_box(slide, 0.8, y, 11, 0.6, f"  {item}", 14, False, RGBColor(0x33, 0x33, 0x33))


def _add_quote_slide(prs, layout, title, quote, bg):
    slide = prs.slides.add_slide(layout)
    _add_bg(slide, bg or "1A3C6E")
    from pptx.enum.text import PP_ALIGN
    _add_text_box(slide, 1, 2, 11, 3, f'"{quote}"', 28, False, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.CENTER)
    if title:
        _add_text_box(slide, 1, 5.5, 11, 0.5, f"-- {title}", 16, False, RGBColor(0xAA, 0xCC, 0xEE), PP_ALIGN.RIGHT)


def _add_thankyou_slide(prs, layout, title, content):
    slide = prs.slides.add_slide(layout)
    _add_bg(slide, "1A3C6E")
    from pptx.enum.text import PP_ALIGN
    _add_text_box(slide, 1, 2, 11, 1.5, title or "Gracias", 48, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.CENTER)
    if content:
        _add_text_box(slide, 1, 4.5, 11, 1, content, 18, False, RGBColor(0xAA, 0xCC, 0xEE), PP_ALIGN.CENTER)


def _resolve_output_dir(path_raw: str) -> Path:
    if path_raw:
        p = Path(path_raw)
        if not p.is_absolute():
            try:
                from actions.path_helper import get_desktop_path
                desk = Path(get_desktop_path())
            except Exception:
                desk = Path.home() / "Desktop"
            p = desk / path_raw
        return p
    try:
        from actions.path_helper import get_desktop_path
        desk = Path(get_desktop_path())
    except Exception:
        desk = Path.home() / "Desktop"
    return desk / "ERIS_Presentaciones"


def _list_templates() -> str:
    return (
        "Acciones de presentation_generator:\n\n"
        "1. create: Crear presentacion PowerPoint\n"
        "   title (str): Titulo de la presentacion\n"
        "   subtitle (str): Subtitulo opcional\n"
        "   author (str): Autor (default: ERIS)\n"
        "   slides (str): JSON con diapositivas:\n"
        '     [{"type":"bullets|section|timeline|quote|thankyou|numbered",\n'
        '       "title":"Titulo","items":["item1","item2"],\n'
        '       "content":"texto","bg_color":"1A3C6E"}\n'
        "     ]\n"
        "   filename (str): Nombre del archivo\n"
        "   output_path (str): Carpeta destino (default: Desktop\\ERIS_Presentaciones)\n\n"
        "Tipos de diapositiva:\n"
        "  title -> portada automatica\n"
        "  section -> separador con fondo azul\n"
        "  bullets -> lista de viñetas\n"
        "  numbered -> lista numerada\n"
        "  timeline -> linea de tiempo (items: {year, desc})\n"
        "  quote -> cita con fondo azul\n"
        "  thankyou -> despedida\n\n"
        "Ejemplo:\n"
        '  presentation_generator action=create title="Mi Tema" slides=\'[{"type":"bullets","title":"Puntos","items":["A","B"]}]\''
    )
