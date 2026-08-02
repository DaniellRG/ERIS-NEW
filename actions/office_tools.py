# -*- coding: utf-8 -*-
"""
office_tools.py — ERIS Excel / Word / PowerPoint por voz.

Acciones:
  excel_create  — Crear un archivo XLSX con datos (params: path, headers, rows)
  excel_read    — Leer celdas de un XLSX (params: path, sheet, max_rows)
  excel_write   — Escribir un valor o añadir fila (params: path, sheet, cell, value | row)
  word_create   — Crear un DOCX (params: path, title, paragraphs)
  word_read     — Leer texto de un DOCX (params: path)
  pptx_create   — Crear un PPTX (params: path, title, slides)
  help          — Listar acciones disponibles
"""
import json
from datetime import datetime
from pathlib import Path

BASE_DIR = Path(__file__).resolve().parent.parent
DOCS_DIR = BASE_DIR / "data" / "documents"


def _resolve_path(path: str) -> Path:
    """Resolve a path, defaulting to data/documents for relative ones."""
    if not path:
        raise ValueError("Se requiere 'path'.")
    p = Path(path).expanduser()
    if not p.is_absolute():
        p = DOCS_DIR / p
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def _as_list(value) -> list:
    """Accept a real list, a JSON-encoded string, or comma/newline separated text."""
    if value is None:
        return []
    if isinstance(value, list):
        return value
    if isinstance(value, str):
        text = value.strip()
        if not text:
            return []
        try:
            parsed = json.loads(text)
            if isinstance(parsed, list):
                return parsed
        except Exception:
            pass
        for sep in ("\n", "|", ","):
            if sep in text:
                parts = [part.strip().strip("[]") for part in text.split(sep) if part.strip()]
                return parts
        return [text]
    return list(value)


def _parse_rows(value):
    """Parse filas: lista real, JSON, o arrays anidados sin comillas ('[[Enero, 100], [Febrero, 150]]')."""
    import re as _re

    if value is None:
        return None
    if isinstance(value, list):
        return value
    text = str(value).strip()
    if not text:
        return None
    try:
        parsed = json.loads(text)
        if isinstance(parsed, list):
            return parsed
    except Exception:
        pass
    if text.startswith("[[") and text.endswith("]]"):
        inner = text[1:-1]
        parts = _re.split(r"\]\s*,\s*\[", inner)
        rows = []
        for part in parts:
            part = part.strip().strip("[]")
            items = [c.strip().strip("\"'") for c in part.split(",") if c.strip()]
            rows.append(items)
        if rows:
            return rows
    return None


def _excel_create(params: dict) -> str:
    path = _resolve_path(params.get("path", ""))
    if path.suffix.lower() != ".xlsx":
        path = path.with_suffix(".xlsx")
    sheet_name = params.get("sheet", "Hoja1")
    headers = _as_list(params.get("headers"))
    rows = _parse_rows(params.get("rows"))
    if rows is None:
        rows = _parse_rows(params.get("data"))
    if rows is None:
        rows = _as_list(params.get("rows"))
    for row in rows:
        if not isinstance(row, list):
            rows = [_as_list(r) for r in rows]
            break

    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet_name[:31]
    if headers:
        ws.append(headers)
    for row in rows:
        ws.append(row)
    wb.save(str(path))

    n_cells = (len(headers) if headers else 0) + len(rows)
    return (
        f"Excel creado: {path}\n"
        f"  Hoja: {sheet_name}\n"
        f"  Columnas: {len(headers)}\n"
        f"  Filas de datos: {len(rows)}\n"
        f"  Celdas escritas: {n_cells}"
    )


def _excel_read(params: dict) -> str:
    path = _resolve_path(params.get("path", ""))
    if not path.exists():
        return f"Archivo no encontrado: {path}"
    sheet_name = params.get("sheet", "")
    max_rows = int(params.get("max_rows", 20))

    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    if sheet_name:
        if sheet_name not in wb.sheetnames:
            return f"Hoja '{sheet_name}' no existe. Hojas: {', '.join(wb.sheetnames)}"
        sheets = [wb[sheet_name]]
    else:
        sheets = wb.worksheets

    lines = [f"Excel: {path}"]
    for ws in sheets:
        lines.append(f"  Hoja '{ws.title}': {ws.max_row} filas x {ws.max_column} columnas")
        count = 0
        for row in ws.iter_rows(min_row=1, max_row=min(max_rows, ws.max_row or 0), values_only=True):
            vals = ["." if v is None else str(v) for v in row]
            lines.append(f"    F{count + 1}: " + " | ".join(vals))
            count += 1
        if (ws.max_row or 0) > max_rows:
            lines.append(f"    ... ({ws.max_row - max_rows} filas más)")
    wb.close()
    return "\n".join(lines)


def _excel_write(params: dict) -> str:
    path = _resolve_path(params.get("path", ""))
    if path.suffix.lower() != ".xlsx":
        path = path.with_suffix(".xlsx")
    sheet_name = params.get("sheet", "Hoja1")
    cell = params.get("cell", "")
    value = params.get("value")
    row = _as_list(params.get("row"))

    import openpyxl
    if path.exists():
        wb = openpyxl.load_workbook(str(path))
        if sheet_name not in wb.sheetnames:
            ws = wb.create_sheet(sheet_name[:31])
        else:
            ws = wb[sheet_name]
    else:
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = sheet_name[:31]

    written = []
    if cell:
        ws[cell] = value
        written.append(f"{cell}={value}")
    if row:
        if not isinstance(row, list):
            return "Error: 'row' debe ser una lista de valores."
        ws.append(row)
        written.append(f"fila en la fila {ws.max_row}")

    if not written:
        return "Error: necesitas 'cell'+'value' o 'row'."
    wb.save(str(path))
    return f"Excel actualizado: {path}\n  Cambios: {', '.join(written)}"


def _word_create(params: dict) -> str:
    path = _resolve_path(params.get("path", ""))
    if path.suffix.lower() != ".docx":
        path = path.with_suffix(".docx")
    title = params.get("title", "")
    paragraphs = _as_list(params.get("paragraphs"))

    from docx import Document
    doc = Document()
    if title:
        doc.add_heading(title, level=0)
    for p in paragraphs:
        doc.add_paragraph(p)
    doc.save(str(path))
    return f"Documento Word creado: {path}\n  Título: {title or '(sin título)'}\n  Párrafos: {len(paragraphs)}"


def _word_read(params: dict) -> str:
    path = _resolve_path(params.get("path", ""))
    if not path.exists():
        return f"Archivo no encontrado: {path}"

    from docx import Document
    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    lines = [f"Documento Word: {path} ({len(paragraphs)} párrafos)"]
    max_p = int(params.get("max_paragraphs", 25))
    for p in paragraphs[:max_p]:
        lines.append(f"  - {p[:200]}")
    if len(paragraphs) > max_p:
        lines.append(f"  ... ({len(paragraphs) - max_p} párrafos más)")
    return "\n".join(lines)


def _pptx_create(params: dict) -> str:
    path = _resolve_path(params.get("path", ""))
    if path.suffix.lower() != ".pptx":
        path = path.with_suffix(".pptx")
    title = params.get("title", "")
    slides = _as_list(params.get("slides"))
    if all(isinstance(s, str) for s in slides):
        try:
            slides = json.loads("[" + ", ".join(slides) + "]")
        except Exception:
            pass
    slides = [s for s in slides if isinstance(s, dict)]

    from pptx import Presentation
    prs = Presentation()
    if title:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = s.get("title", "")
        bullets = s.get("bullets", []) or []
        tf = slide.placeholders[1].text_frame
        tf.text = ""
        for i, b in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = b
    prs.save(str(path))
    return f"Presentación creada: {path}\n  Título: {title or '(sin título)'}\n  Diapositivas: {len(slides)}"


def office_docs(parameters: dict, player=None) -> str:
    """
    Excel / Word / PowerPoint por voz.
    Acciones: excel_create, excel_read, excel_write, word_create, word_read, pptx_create, help
    """
    action = str(parameters.get("action", "help")).lower()

    if action == "help":
        return (
            "Office Docs — acciones disponibles:\n"
            "  excel_create: path, sheet, headers (lista), rows (lista de filas)\n"
            "  excel_read: path, sheet, max_rows\n"
            "  excel_write: path, sheet, cell (ej 'B2'), value | row (lista)\n"
            "  word_create: path, title, paragraphs (lista)\n"
            "  word_read: path, max_paragraphs\n"
            "  pptx_create: path, title, slides (lista de {title, bullets})\n"
            "Ejemplo: {action: excel_create, path: gastos, headers: [Mes, Total], rows: [[Enero, 100], [Febrero, 150]]}"
        )

    handlers = {
        "excel_create": _excel_create,
        "excel_read": _excel_read,
        "excel_write": _excel_write,
        "word_create": _word_create,
        "word_read": _word_read,
        "pptx_create": _pptx_create,
    }

    handler = handlers.get(action)
    if not handler:
        return f"Acción desconocida: '{action}'. Usá 'help' para ver las acciones."

    try:
        if player:
            player.write_log(f"[office_docs] {action} ...")
        result = handler(parameters)
        if player:
            player.write_log(f"[office_docs] {action} OK")
        return result
    except ImportError as e:
        return f"Falta la librería: {e}. Instalá openpyxl, python-docx o python-pptx."
    except Exception as e:
        return f"Error en {action}: {e}"
